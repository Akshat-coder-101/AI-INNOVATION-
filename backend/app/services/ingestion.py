import io
import os
import uuid
import re
import logging
from typing import List, Dict, Any, Tuple, Optional, Sequence, Union, TypeVar
from pypdf import PdfReader
import docx
import pptx
from sqlalchemy.orm import Session
from fastapi import HTTPException, status
from ..database import DBMaterial, DBMaterialChunk
from ..config import settings
from .llm import LLMService

TChunk = TypeVar("TChunk", bound=Union[DBMaterialChunk, Dict[str, Any]])

logger = logging.getLogger("sahayak.ingestion")

ALLOWED_EXTENSIONS = {"pdf", "docx", "pptx", "txt", "doc", "ppt", "md", "markdown"}

def clean_text(text: str) -> str:
    text = re.sub(r'\s+', ' ', text)
    return text.strip()

def chunk_text(text: str, chunk_size_words: int = 250, overlap_words: int = 40) -> List[str]:
    words = text.split()
    chunks = []
    if not words:
        return chunks
    
    start = 0
    while start < len(words):
        end = min(start + chunk_size_words, len(words))
        chunk = " ".join(words[start:end])
        if len(chunk.strip()) > 30:
            chunks.append(chunk)
        if end >= len(words):
            break
        start += chunk_size_words - overlap_words
    return chunks

class IngestionService:
    @staticmethod
    def validate_file(filename: str, content: bytes) -> str:
        """Validates file presence, size limit, and supported extension."""
        if not content or len(content) == 0:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Uploaded file is empty."
            )
        
        max_bytes = settings.MAX_UPLOAD_MB * 1024 * 1024
        if len(content) > max_bytes:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"File exceeds maximum allowed size of {settings.MAX_UPLOAD_MB}MB."
            )
        
        parts = filename.rsplit(".", 1)
        if len(parts) < 2:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="File has no extension. Allowed extensions: PDF, DOCX, PPTX, TXT, MD, DOC, PPT."
            )
        
        ext = parts[1].lower().strip()
        if ext not in ALLOWED_EXTENSIONS:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Unsupported file extension '.{ext}'. Allowed extensions: PDF, DOCX, PPTX, TXT, MD, DOC, PPT."
            )
        return ext

    @staticmethod
    def save_file_safely(filename: str, content: bytes, ext: str) -> str:
        """Saves file to DOC_STORAGE_DIR using a secure, traversal-safe random UUID name."""
        backend_dir = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
        storage_dir = os.path.join(backend_dir, settings.DOC_STORAGE_DIR)
        os.makedirs(storage_dir, exist_ok=True)
        safe_filename = f"{uuid.uuid4().hex}.{ext}"
        target_path = os.path.join(storage_dir, safe_filename)
        try:
            with open(target_path, "wb") as f:
                f.write(content)
            return target_path
        except Exception as e:
            logger.warning(f"Could not persist upload to disk ({e}); proceeding in-memory.")
            return ""

    @staticmethod
    def parse_pdf(file_bytes: bytes) -> List[Dict[str, Any]]:
        reader = PdfReader(io.BytesIO(file_bytes))
        pages_content = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():
                pages_content.append({
                    "page_number": i + 1,
                    "title": f"Chapter/Page {i + 1}",
                    "text": clean_text(text),
                    "content_type": "page"
                })
        return pages_content

    @staticmethod
    def parse_docx(file_bytes: bytes) -> List[Dict[str, Any]]:
        """
        Extracts paragraphs, headings, bullet/numbered lists, and tables from DOCX,
        preserving structural hierarchy and document order.
        """
        doc = docx.Document(io.BytesIO(file_bytes))
        sections: List[Dict[str, Any]] = []
        current_heading = "Document Content"
        current_blocks: List[str] = []
        
        # Iterate over body elements in sequential order
        from docx.text.paragraph import Paragraph
        from docx.table import Table

        for child in doc._element.body:
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if tag == "p":
                para = Paragraph(child, doc)
                text = para.text.strip()
                if not text:
                    continue
                style_name = (para.style.name or "").lower() if para.style else ""
                
                # Detect headings
                if "heading" in style_name or "title" in style_name:
                    if current_blocks:
                        sections.append({
                            "heading": current_heading,
                            "text": "\n\n".join(current_blocks),
                            "content_type": "section"
                        })
                        current_blocks = []
                    current_heading = text
                else:
                    # Detect list items
                    if "list" in style_name or "bullet" in style_name:
                        current_blocks.append(f"• {text}")
                    else:
                        current_blocks.append(text)
                        
            elif tag == "tbl":
                table = Table(child, doc)
                table_rows = []
                for row in table.rows:
                    row_cells = [clean_text(cell.text) for cell in row.cells]
                    if any(row_cells):
                        table_rows.append(row_cells)
                
                if table_rows:
                    # Format as Markdown Table
                    header = table_rows[0]
                    md_table_lines = ["| " + " | ".join(header) + " |"]
                    md_table_lines.append("| " + " | ".join(["---"] * len(header)) + " |")
                    for row in table_rows[1:]:
                        # Ensure cell count matches header
                        padded_row = row + [""] * (len(header) - len(row))
                        md_table_lines.append("| " + " | ".join(padded_row[:len(header)]) + " |")
                    
                    current_blocks.append("\n".join(md_table_lines))

        if current_blocks:
            sections.append({
                "heading": current_heading,
                "text": "\n\n".join(current_blocks),
                "content_type": "section"
            })
            
        if not sections:
            # Fallback if no structured elements extracted
            raw_text = "\n".join([p.text.strip() for p in doc.paragraphs if p.text.strip()])
            if raw_text:
                sections.append({
                    "heading": "Document Content",
                    "text": clean_text(raw_text),
                    "content_type": "text"
                })

        return sections

    @staticmethod
    def parse_pptx(file_bytes: bytes) -> List[Dict[str, Any]]:
        """
        Extracts slide titles, bullet points, tables, and speaker notes from PPTX.
        """
        prs = pptx.Presentation(io.BytesIO(file_bytes))
        slides: List[Dict[str, Any]] = []
        
        for i, slide in enumerate(prs.slides):
            slide_number = i + 1
            title = f"Slide {slide_number}"
            slide_blocks: List[str] = []
            
            # 1. Slide Title
            title_shape = getattr(slide.shapes, "title", None)
            if title_shape and getattr(title_shape, "text", None) and title_shape.text.strip():
                title = clean_text(title_shape.text.strip())
            else:
                for s_shape in slide.shapes:
                    if getattr(s_shape, "has_text_frame", False):
                        tf = getattr(s_shape, "text_frame", None)
                        if tf and getattr(tf, "text", None) and tf.text.strip():
                            first_line = tf.text.strip().split("\n")[0].strip()
                            if len(first_line) > 2:
                                title = clean_text(first_line)
                                break
            
            # 2. Extract Text Frames & Tables
            for shape in slide.shapes:
                # Shape is title shape already handled? Don't duplicate unless it has extra info
                if shape == slide.shapes.title:
                    continue
                    
                if getattr(shape, "has_text_frame", False):
                    tf = getattr(shape, "text_frame", None)
                    if tf:
                        for paragraph in tf.paragraphs:
                            p_text = paragraph.text.strip()
                            if p_text:
                                indent_prefix = "  " * getattr(paragraph, "level", 0)
                                bullet_char = "• " if getattr(paragraph, "level", 0) > 0 else ""
                                slide_blocks.append(f"{indent_prefix}{bullet_char}{p_text}")
                                
                elif getattr(shape, "has_table", False):
                    tbl = getattr(shape, "table", None)
                    if tbl and tbl.rows:
                        table_rows = []
                        for row in tbl.rows:
                            row_cells = [clean_text(cell.text) for cell in row.cells]
                            if any(row_cells):
                                table_rows.append(row_cells)
                        if table_rows:
                            header = table_rows[0]
                            md_table = ["| " + " | ".join(header) + " |"]
                            md_table.append("| " + " | ".join(["---"] * len(header)) + " |")
                            for r in table_rows[1:]:
                                padded = r + [""] * (len(header) - len(r))
                                md_table.append("| " + " | ".join(padded[:len(header)]) + " |")
                            slide_blocks.append("\n".join(md_table))

            # 3. Extract Speaker Notes (if available)
            if getattr(slide, "has_notes_slide", False) and slide.notes_slide:
                notes_tf = getattr(slide.notes_slide, "notes_text_frame", None)
                if notes_tf and notes_tf.text.strip():
                    notes_content = clean_text(notes_tf.text.strip())
                    slide_blocks.append(f"[Speaker Notes]: {notes_content}")

            content_text = "\n\n".join(slide_blocks)
            slides.append({
                "slide_number": slide_number,
                "title": title,
                "text": content_text if content_text else title,
                "content_type": "slide"
            })
            
        return slides

    @staticmethod
    def parse_markdown_or_txt(file_bytes: bytes) -> List[Dict[str, Any]]:
        """
        Extracts structured sections from Markdown or plain text files by header tokens (#, ##, ###).
        """
        text = file_bytes.decode("utf-8", errors="ignore")
        lines = text.splitlines()
        sections: List[Dict[str, Any]] = []
        current_heading = "Document Content"
        current_lines: List[str] = []

        for line in lines:
            header_match = re.match(r'^(#{1,4})\s+(.+)$', line.strip())
            if header_match:
                if current_lines:
                    sections.append({
                        "heading": current_heading,
                        "text": clean_text("\n".join(current_lines)),
                        "content_type": "section"
                    })
                    current_lines = []
                current_heading = header_match.group(2).strip()
            else:
                if line.strip():
                    current_lines.append(line.strip())

        if current_lines:
            sections.append({
                "heading": current_heading,
                "text": clean_text("\n".join(current_lines)),
                "content_type": "section"
            })

        if not sections:
            sections.append({
                "heading": "Document Content",
                "text": clean_text(text),
                "content_type": "text"
            })
            
        return sections

    @classmethod
    async def extract_title_and_topics(
        cls, 
        filename: str, 
        preview_text: str, 
        chapters: List[Dict[str, Any]]
    ) -> Tuple[str, List[str]]:
        """Uses a guarded LLM pass to extract detected title and key topics, with deterministic fallback."""
        # 1. Fallback heuristic defaults
        base_name = os.path.splitext(filename)[0].replace("_", " ").replace("-", " ").title()
        fallback_title = chapters[0]["title"] if chapters and "Chapter" not in chapters[0]["title"] else base_name
        fallback_topics = [ch["title"] for ch in chapters[:5] if "title" in ch and "Chapter" not in ch["title"]]
        if not fallback_topics:
            words = re.findall(r'\b[A-Z][a-z]{3,}\b', preview_text)
            fallback_topics = list(dict.fromkeys(words))[:4] if words else [base_name, "Key Concepts", "Core Principles"]

        # 2. Guarded LLM Extraction
        try:
            system_prompt = (
                "You are an expert document indexer. Analyze the document excerpt and extract a concise document title "
                "and 3 to 6 key educational topics covered. Return ONLY valid JSON: {\"detected_title\": \"...\", \"key_topics\": [\"...\"]}"
            )
            user_prompt = f"Filename: {filename}\n\nDocument Excerpt:\n{preview_text[:2000]}"
            res = await LLMService.generate_json(
                system_prompt=system_prompt,
                user_prompt=user_prompt,
                schema_hint='{"detected_title": "...", "key_topics": ["Topic 1", "Topic 2"]}',
                temperature=0.2
            )
            if res and isinstance(res, dict):
                title = res.get("detected_title") or fallback_title
                topics = res.get("key_topics") or fallback_topics
                if isinstance(topics, list) and len(topics) > 0:
                    return str(title).strip(), [str(t).strip() for t in topics if str(t).strip()]
        except Exception as e:
            logger.info(f"LLM title/topic extraction skipped/fallback: {e}")

        return fallback_title, fallback_topics

    @classmethod
    def _safe_convert_legacy_binary(cls, content: bytes, ext: str) -> bytes:
        """
        Attempts headless conversion for legacy .doc / .ppt if LibreOffice / soffice is available.
        Otherwise extracts text tokens safely or raises informative error.
        """
        import shutil
        import tempfile
        import subprocess

        soffice_bin = shutil.which("soffice") or shutil.which("libreoffice")
        if soffice_bin:
            with tempfile.TemporaryDirectory() as tmpdir:
                in_path = os.path.join(tmpdir, f"input.{ext}")
                with open(in_path, "wb") as f:
                    f.write(content)
                target_fmt = "docx" if ext == "doc" else "pptx"
                cmd = [soffice_bin, "--headless", "--convert-to", target_fmt, "--outdir", tmpdir, in_path]
                res = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=20)
                out_path = os.path.join(tmpdir, f"input.{target_fmt}")
                if res.returncode == 0 and os.path.exists(out_path):
                    with open(out_path, "rb") as f:
                        return f.read()

        # Fallback: check if content is readable text or UTF-8/Latin-1
        text_matches = re.findall(rb'[\x20-\x7E\t\n\r]{4,}', content)
        if text_matches and len(text_matches) > 5:
            extracted_text = " ".join([m.decode("latin-1", errors="ignore") for m in text_matches])
            return extracted_text.encode("utf-8")

        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Legacy binary .{ext} files require conversion to modern OpenXML (.{ext}x) or PDF format."
        )

    @classmethod
    def process_file(cls, filename: str, content: bytes, db: Session) -> Dict[str, Any]:
        material_id = str(uuid.uuid4())
        ext = filename.split(".")[-1].lower() if "." in filename else "txt"
        
        raw_chunks_with_meta: List[Dict[str, Any]] = []
        detected_chapters: List[Dict[str, Any]] = []
        all_text = ""
        
        if ext == "pdf":
            pages = cls.parse_pdf(content)
            all_text = "\n".join([str(p.get("text", "")) for p in pages])
            for p in pages:
                sub_chunks = chunk_text(p["text"])
                chapter_name = p.get("title") or f"Chapter/Page {p['page_number']}"
                detected_chapters.append({
                    "title": chapter_name,
                    "page": p["page_number"],
                    "preview": p["text"][:150] + "..."
                })
                for idx, chunk in enumerate(sub_chunks):
                    raw_chunks_with_meta.append({
                        "chapter": chapter_name,
                        "page": p["page_number"],
                        "section": f"Section {idx+1}",
                        "content": chunk,
                        "content_type": "page"
                    })
                    
        elif ext in ["docx", "doc"]:
            try:
                sections = cls.parse_docx(content)
            except Exception as doc_err:
                logger.warning(f"DOCX parse failed ({doc_err}), attempting legacy converter/fallback.")
                converted_bytes = cls._safe_convert_legacy_binary(content, ext)
                try:
                    sections = cls.parse_docx(converted_bytes)
                except Exception:
                    sections = cls.parse_markdown_or_txt(converted_bytes)

            all_text = "\n\n".join([str(s.get("text", "")) for s in sections])
            for i, s in enumerate(sections):
                heading_title = s["heading"]
                detected_chapters.append({
                    "title": heading_title,
                    "page": i + 1,
                    "preview": s["text"][:150] + "..."
                })
                sub_chunks = chunk_text(s["text"])
                for idx, chunk in enumerate(sub_chunks):
                    raw_chunks_with_meta.append({
                        "chapter": heading_title,
                        "page": i + 1,
                        "section": f"Section {idx+1}",
                        "content": chunk,
                        "content_type": s.get("content_type", "section")
                    })
                    
        elif ext in ["pptx", "ppt"]:
            slides: List[Dict[str, Any]] = []
            try:
                slides = cls.parse_pptx(content)
            except Exception as ppt_err:
                logger.warning(f"PPTX parse failed ({ppt_err}), attempting legacy converter/fallback.")
                converted_bytes = cls._safe_convert_legacy_binary(content, ext)
                try:
                    slides = cls.parse_pptx(converted_bytes)
                except Exception:
                    raw_slides = cls.parse_markdown_or_txt(converted_bytes)
                    slides = [{"slide_number": i + 1, "title": s["heading"], "text": s["text"], "content_type": "slide"} for i, s in enumerate(raw_slides)]

            all_text = "\n\n".join([str(s.get("text", "")) for s in slides])
            for s in slides:
                slide_title = str(s.get("title", ""))
                slide_text = str(s.get("text", ""))
                slide_num = s.get("slide_number", 1)
                detected_chapters.append({
                    "title": slide_title,
                    "page": slide_num,
                    "preview": slide_text[:150] + "..."
                })
                sub_chunks = chunk_text(slide_text)
                for idx, chunk in enumerate(sub_chunks):
                    raw_chunks_with_meta.append({
                        "chapter": slide_title,
                        "page": slide_num,
                        "section": f"Slide {slide_num}",
                        "content": chunk,
                        "content_type": "slide"
                    })

        elif ext in ["md", "markdown"]:
            sections = cls.parse_markdown_or_txt(content)
            all_text = "\n\n".join([str(s.get("text", "")) for s in sections])
            for i, s in enumerate(sections):
                heading_title = s["heading"]
                detected_chapters.append({
                    "title": heading_title,
                    "page": i + 1,
                    "preview": s["text"][:150] + "..."
                })
                sub_chunks = chunk_text(s["text"])
                for idx, chunk in enumerate(sub_chunks):
                    raw_chunks_with_meta.append({
                        "chapter": heading_title,
                        "page": i + 1,
                        "section": f"Section {idx+1}",
                        "content": chunk,
                        "content_type": "markdown"
                    })
                    
        else: # Plain text or generic fallback
            sections = cls.parse_markdown_or_txt(content)
            all_text = "\n\n".join([str(s.get("text", "")) for s in sections])
            for i, s in enumerate(sections):
                heading_title = s["heading"]
                detected_chapters.append({
                    "title": heading_title,
                    "page": i + 1,
                    "preview": s["text"][:150] + "..."
                })
                sub_chunks = chunk_text(s["text"])
                for idx, chunk in enumerate(sub_chunks):
                    raw_chunks_with_meta.append({
                        "chapter": heading_title,
                        "page": i + 1,
                        "section": f"Section {idx+1}",
                        "content": chunk,
                        "content_type": "text"
                    })

        # Save material record
        db_material = DBMaterial(
            id=material_id,
            filename=filename,
            content_type=ext,
            total_sections=len(detected_chapters),
            raw_text=all_text[:50000] # Store preview
        )
        db.add(db_material)

        # Save chunks with semantic vector representation
        from .rag import RAGService
        for item in raw_chunks_with_meta:
            embedding = RAGService.generate_embedding(item["content"])
            db_chunk = DBMaterialChunk(
                id=str(uuid.uuid4()),
                material_id=material_id,
                chapter=item["chapter"],
                page=item.get("page", 1),
                section=item.get("section", ""),
                content=item["content"],
                embedding=embedding,
                token_count=len(item["content"].split())
            )
            db.add(db_chunk)

        db.commit()

        return {
            "material_id": material_id,
            "filename": filename,
            "total_pages_or_sections": len(detected_chapters),
            "chunks_count": len(raw_chunks_with_meta),
            "chapters": detected_chapters[:20],
            "preview": all_text[:300] + "..."
        }

    @classmethod
    def filter_chunks_by_chapter_or_topic(
        cls, 
        chunks: Sequence[TChunk], 
        target_chapter: Optional[str]
    ) -> List[TChunk]:
        """
        Filters chunks when a specific chapter or section is requested by the student (e.g. 'Chapter 4').
        Extracts chapter number or keyword and matches against chunk.chapter, chunk.section, or content.
        Falls back gracefully to all chunks if no specific match is found.
        """
        if not target_chapter or not target_chapter.strip() or not chunks:
            return list(chunks)

        target_norm = target_chapter.strip().lower()
        
        # Check for chapter number (e.g. "4" from "Chapter 4" or "ch 4" or "unit 4")
        num_match = re.search(r'\b(?:chapter|ch|unit|section|part|page)?\s*(\d+)\b', target_norm)
        ch_num = num_match.group(1) if num_match else None

        matched_chunks = []

        def _get_field(chunk_item, field_name, default=""):
            if hasattr(chunk_item, field_name):
                return getattr(chunk_item, field_name) or default
            if isinstance(chunk_item, dict):
                if field_name in chunk_item:
                    return chunk_item[field_name] or default
                meta = chunk_item.get("metadata")
                if isinstance(meta, dict) and field_name in meta:
                    return meta[field_name] or default
            return default

        for c in chunks:
            ch_str = str(_get_field(c, "chapter", "")).lower()
            sec_str = str(_get_field(c, "section", "")).lower()
            page_val = str(_get_field(c, "page", ""))
            content_str = str(_get_field(c, "content", "")).lower()
            
            # Exact or substring match in chapter title or section
            if target_norm in ch_str or target_norm in sec_str or target_norm in content_str:
                matched_chunks.append(c)
                continue
                
            # Number match (e.g. page or chapter number)
            if ch_num:
                if (f"chapter {ch_num}" in ch_str or f"chapter {ch_num}" in content_str or 
                    f"page {ch_num}" in ch_str or f"slide {ch_num}" in ch_str or 
                    f"part {ch_num}" in ch_str or f"section {ch_num}" in ch_str or
                    ch_str == ch_num):
                    matched_chunks.append(c)
                    continue
                if page_val == ch_num:
                    matched_chunks.append(c)
                    continue

        if matched_chunks:
            logger.info(f"[IngestionService] Filtered {len(matched_chunks)}/{len(chunks)} chunks for target '{target_chapter}'.")
            return matched_chunks

        # Keyword match fallback
        target_keywords = [w for w in re.findall(r'\w+', target_norm) if len(w) > 3 and w not in ["chapter", "section", "part", "unit", "teach", "learn"]]
        if target_keywords:
            kw_matches = [
                c for c in chunks 
                if any(k in str(_get_field(c, "content", "")).lower() or k in str(_get_field(c, "chapter", "")).lower() for k in target_keywords)
            ]
            if kw_matches:
                logger.info(f"[IngestionService] Keyword-matched {len(kw_matches)}/{len(chunks)} chunks for target '{target_chapter}'.")
                return kw_matches

        return list(chunks)

    @classmethod
    async def process_document_upload(cls, filename: str, content: bytes, db: Session) -> Dict[str, Any]:
        """Validates upload, stores securely, chunks, embeds, extracts title/topics, and returns canonical response."""
        ext = cls.validate_file(filename, content)
        cls.save_file_safely(filename, content, ext)
        
        ingest_res = cls.process_file(filename, content, db)
        doc_id = ingest_res["material_id"]
        
        detected_title, key_topics = await cls.extract_title_and_topics(
            filename=filename,
            preview_text=ingest_res.get("preview", ""),
            chapters=ingest_res.get("chapters", [])
        )
        
        return {
            "document_id": doc_id,
            "filename": filename,
            "page_count": max(1, ingest_res.get("total_pages_or_sections", 1)),
            "chunk_count": ingest_res.get("chunks_count", 0),
            "detected_title": detected_title,
            "key_topics": key_topics
        }


