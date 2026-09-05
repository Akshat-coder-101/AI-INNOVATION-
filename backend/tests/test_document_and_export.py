import io
import os
import sys
import uuid
import pytest
import docx
import pptx
from fastapi.testclient import TestClient

# Add backend directory to sys.path
sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), "..")))

from main import app
from app.database import SessionLocal, DBMaterial, DBMaterialChunk, DBLessonSession, DBExportJob, init_db
from app.services.ingestion import IngestionService
from app.services.video import VideoService
from app.state_machine.teacher_agent import TeacherAgentStateMachine

client = TestClient(app)

@pytest.fixture(autouse=True)
def ensure_db():
    init_db()

# --- Helper functions to create in-memory test documents ---

def create_sample_docx() -> bytes:
    doc = docx.Document()
    doc.add_heading("Chapter 1: Neural Networks and Deep Learning", level=1)
    doc.add_paragraph("Neural networks are computational models inspired by biological neural circuits.")
    
    # Bullet points
    doc.add_paragraph("Perceptron architecture", style="List Bullet")
    doc.add_paragraph("Activation functions like ReLU and Sigmoid", style="List Bullet")
    doc.add_paragraph("Backpropagation and gradient descent", style="List Bullet")
    
    doc.add_heading("Section 2: Performance Evaluation Matrix", level=2)
    doc.add_paragraph("Below is the comparative metric breakdown across network architectures:")
    
    # Table
    table = doc.add_table(rows=3, cols=3)
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = "Architecture"
    hdr_cells[1].text = "Parameters"
    hdr_cells[2].text = "Accuracy"
    
    row1 = table.rows[1].cells
    row1[0].text = "MLP Baseline"
    row1[1].text = "1.2M"
    row1[2].text = "88.4%"
    
    row2 = table.rows[2].cells
    row2[0].text = "Deep ResNet"
    row2[1].text = "25.6M"
    row2[2].text = "96.2%"
    
    doc.add_paragraph("In conclusion, deeper residual networks achieve superior generalization.")
    
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()

def create_sample_pptx() -> bytes:
    prs = pptx.Presentation()
    
    # Slide 1: Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    slide1.shapes.title.text = "Introduction to Quantum Computing"
    slide1.placeholders[1].text = "Foundations of Qubits and Superposition\nBy Prof. Sahayak"
    
    # Slide 2: Bullet slide with speaker notes
    bullet_slide_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    slide2.shapes.title.text = "Qubits vs Classical Bits"
    tf = slide2.shapes.placeholders[1].text_frame
    tf.text = "Classical bits exist strictly in state 0 or 1"
    p = tf.add_paragraph()
    p.text = "Qubits leverage quantum superposition |ψ⟩ = α|0⟩ + β|1⟩"
    p.level = 1
    p2 = tf.add_paragraph()
    p2.text = "Entanglement enables non-local quantum state correlations"
    p2.level = 1
    
    # Add speaker notes to Slide 2
    notes_slide = slide2.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "Emphasize to students that measuring the qubit collapses the wave function."
    
    # Slide 3: Slide with table
    blank_layout = prs.slide_layouts[6]
    slide3 = prs.slides.add_slide(blank_layout)
    
    # Add title box
    txBox = slide3.shapes.add_textbox(0, 0, pptx.util.Inches(8), pptx.util.Inches(1))
    txBox.text_frame.text = "Quantum Algorithms Summary"
    
    # Add table
    table_shape = slide3.shapes.add_table(3, 2, pptx.util.Inches(1), pptx.util.Inches(1.5), pptx.util.Inches(6), pptx.util.Inches(2))
    table = table_shape.table
    table.cell(0, 0).text = "Algorithm"
    table.cell(0, 1).text = "Speedup"
    table.cell(1, 0).text = "Shor's Algorithm"
    table.cell(1, 1).text = "Exponential"
    table.cell(2, 0).text = "Grover's Search"
    table.cell(2, 1).text = "Quadratic"
    
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()

def create_sample_markdown() -> bytes:
    md_content = """# Unit 1: Principles of Thermodynamics
Thermodynamics governs energy transformations and macroscopic properties of systems.

## Chapter 1: First Law and Conservation of Energy
Energy cannot be created or destroyed; it only transforms from one form to another:
ΔU = Q - W

- **Internal Energy (U)**: Total kinetic and potential energy of molecules
- **Heat (Q)**: Energy transferred due to temperature difference
- **Work (W)**: Energy transferred by mechanical force

## Chapter 2: Second Law and Entropy
Entropy of an isolated system always increases over time in spontaneous processes:
dS >= dQ / T
"""
    return md_content.encode("utf-8")


# --- Test Cases ---

def test_docx_ingestion_and_table_extraction():
    db = SessionLocal()
    try:
        docx_bytes = create_sample_docx()
        sections = IngestionService.parse_docx(docx_bytes)
        
        assert len(sections) >= 2
        # Check headings
        headings = [s["heading"] for s in sections]
        assert any("Neural Networks" in h for h in headings)
        assert any("Performance Evaluation" in h for h in headings)
        
        # Check table extraction in markdown table format
        all_text = " ".join([s["text"] for s in sections])
        assert "MLP Baseline" in all_text
        assert "Deep ResNet" in all_text
        assert "88.4%" in all_text
        assert "Backpropagation" in all_text
        
        # Ingest full document through process_file
        res = IngestionService.process_file("deep_learning_guide.docx", docx_bytes, db)
        assert res["material_id"] is not None
        assert res["chunks_count"] > 0
        
        # Verify chunks stored in DB with metadata
        chunks = db.query(DBMaterialChunk).filter(DBMaterialChunk.material_id == res["material_id"]).all()
        assert len(chunks) > 0
        assert any("ResNet" in c.content for c in chunks)
        assert any("Perceptron" in c.content for c in chunks)
    finally:
        db.close()


def test_pptx_ingestion_slides_notes_and_tables():
    db = SessionLocal()
    try:
        pptx_bytes = create_sample_pptx()
        slides = IngestionService.parse_pptx(pptx_bytes)
        
        assert len(slides) == 3
        # Slide 1
        assert "Quantum Computing" in slides[0]["title"]
        # Slide 2 & Speaker Notes
        assert "Qubits" in slides[1]["title"]
        assert "superposition" in slides[1]["text"].lower()
        assert "Speaker Notes" in slides[1]["text"]
        assert "collapses the wave function" in slides[1]["text"]
        
        # Slide 3 & Table
        assert "Algorithms" in slides[2]["title"]
        assert "Shor's Algorithm" in slides[2]["text"]
        assert "Exponential" in slides[2]["text"]
        
        # Ingest through process_file
        res = IngestionService.process_file("quantum_lecture.pptx", pptx_bytes, db)
        assert res["material_id"] is not None
        assert res["total_pages_or_sections"] == 3
        
        chunks = db.query(DBMaterialChunk).filter(DBMaterialChunk.material_id == res["material_id"]).all()
        assert len(chunks) >= 3
        assert any("Shor" in c.content for c in chunks)
    finally:
        db.close()


def test_markdown_and_txt_ingestion():
    db = SessionLocal()
    try:
        md_bytes = create_sample_markdown()
        sections = IngestionService.parse_markdown_or_txt(md_bytes)
        
        assert len(sections) >= 2
        headings = [s["heading"] for s in sections]
        assert any("First Law" in h or "Thermodynamics" in h for h in headings)
        assert any("Second Law" in h or "Entropy" in h for h in headings)
        
        res = IngestionService.process_file("thermodynamics.md", md_bytes, db)
        assert res["chunks_count"] > 0
        chunks = db.query(DBMaterialChunk).filter(DBMaterialChunk.material_id == res["material_id"]).all()
        assert any("Entropy" in c.content for c in chunks)
    finally:
        db.close()


@pytest.mark.asyncio
async def test_ai_teacher_lesson_grounded_in_docx():
    db = SessionLocal()
    try:
        docx_bytes = create_sample_docx()
        ingest_res = IngestionService.process_file("neural_nets.docx", docx_bytes, db)
        material_id = ingest_res["material_id"]
        
        plan = await TeacherAgentStateMachine.generate_lesson_plan(
            topic=None,
            material_id=material_id,
            profile=None,
            time_budget_minutes=20,
            language="en",
            db=db
        )
        
        assert plan.session_id is not None
        assert len(plan.segments) >= 2
        assert plan.material_id == material_id
        
        # Render first segment
        seg = await TeacherAgentStateMachine.render_segment(
            session_id=plan.session_id,
            segment_id=1,
            language="en",
            db=db
        )
        assert seg.concept is not None
        assert seg.spoken_script is not None
        assert len(seg.captions) > 0
    finally:
        db.close()


def test_upload_documents_api_docx_and_pptx():
    # Test upload endpoint with DOCX
    docx_bytes = create_sample_docx()
    files = {"file": ("deep_learning.docx", docx_bytes, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")}
    res = client.post("/api/documents/upload", files=files)
    assert res.status_code == 200
    data = res.json()
    assert data["document_id"] is not None
    assert data["chunk_count"] > 0
    
    # Test upload endpoint with PPTX
    pptx_bytes = create_sample_pptx()
    files_ppt = {"file": ("quantum_slides.pptx", pptx_bytes, "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
    res_ppt = client.post("/api/documents/upload", files=files_ppt)
    assert res_ppt.status_code == 200
    data_ppt = res_ppt.json()
    assert data_ppt["document_id"] is not None
    assert data_ppt["page_count"] == 3


def test_export_lesson_video_endpoints_and_job_lifecycle():
    db = SessionLocal()
    try:
        # Create a test session
        session_id = str(uuid.uuid4())
        db_sess = DBLessonSession(
            id=session_id,
            topic="Binary Search Trees",
            language="en",
            time_budget=20,
            plan_json={
                "session_id": session_id,
                "topic": "Binary Search Trees",
                "segments": [
                    {"id": 1, "concept": "BST Ordering Invariant", "visual_type": "labeled-diagram", "summary": "Left < Node < Right"},
                    {"id": 2, "concept": "Search & Insertion Complexity", "visual_type": "equation/graph", "summary": "Logarithmic time operations"}
                ]
            }
        )
        db.add(db_sess)
        db.commit()

        # 1. Trigger export
        res = client.post(f"/api/lesson/{session_id}/export")
        assert res.status_code == 200
        export_data = res.json()
        job_id = export_data["job_id"]
        assert job_id is not None
        assert export_data["status"] in ["queued", "processing", "completed", "failed"]

        # 2. Check status polling endpoint
        status_res = client.get(f"/api/lesson/export/{job_id}/status")
        assert status_res.status_code == 200
        status_data = status_res.json()
        assert status_data["job_id"] == job_id
        assert "status" in status_data
        assert "progress" in status_data

        # 3. Test non-existent session 404
        bad_res = client.post(f"/api/lesson/non-existent-id/export")
        assert bad_res.status_code == 404

        # 4. Test non-existent job 404
        bad_job = client.get(f"/api/lesson/export/non-existent-job/status")
        assert bad_job.status_code == 404

        # 5. Test download endpoint validation
        dl_res = client.get(f"/api/lesson/export/{job_id}/download")
        # Since local environment has no ffmpeg or job is queued/failed, download safely returns 400 or 404 with diagnostic
        assert dl_res.status_code in [400, 404, 200]
    finally:
        db.close()


@pytest.mark.asyncio
async def test_video_export_worker_graceful_missing_ffmpeg():
    db = SessionLocal()
    try:
        session_id = str(uuid.uuid4())
        job_id = str(uuid.uuid4())
        db_sess = DBLessonSession(
            id=session_id,
            topic="Test Lesson",
            language="en",
            time_budget=10,
            plan_json={"session_id": session_id, "topic": "Test Lesson", "segments": [{"id": 1, "concept": "Intro"}]}
        )
        job = DBExportJob(id=job_id, session_id=session_id, status="queued", progress=0)
        db.add(db_sess)
        db.add(job)
        db.commit()

        # Execute worker directly
        await VideoService.export_full_lesson_video(job_id=job_id, session_id=session_id)

        # Inspect job record
        updated_job = db.query(DBExportJob).filter(DBExportJob.id == job_id).first()
        assert updated_job is not None
        # On host without ffmpeg, status becomes failed with descriptive message rather than crashing
        assert updated_job.status in ["completed", "failed"]
        if updated_job.status == "failed":
            assert "ffmpeg" in updated_job.error_message.lower() or "not found" in updated_job.error_message.lower()
    finally:
        db.close()
